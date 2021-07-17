from unittest import TestCase

from colour import Color
from ddt import ddt, data, unpack
from pptx import Presentation
from pptx.util import Cm

from source.shape_formatting import ShapeFormatting
from source.activity_layout_attributes import ActivityLayoutAttributes
from source.plan_activity import PlanActivity
from source.plot_driver import PlotDriver
from source.tests.testing_utilities import parse_date
from source.visual_element_shape import VisualElementShape

plan_visual_config_test_data = {
    'vis_cfg_01': {
        'top': Cm(0),
        'left': Cm(0),
        'bottom': Cm(20),
        'right': Cm(30),
        'track_height': Cm(1),
        'track_gap': Cm(0.5),
        'min_start_date': parse_date('2021-01-01'),
        'max_end_date': parse_date('2021-12-31'),
        'milestone_width': Cm(0.4),
        'milestone_text_width': Cm(0.5),
        'activity_text_width': Cm(5),
        'text_margin': Cm(0.2),
        'activity_shape': 'RECTANGLE',
        'milestone_shape': 'DIAMOND'
    }
}
display_attribute_test_data = {
    'disp_01': {
        'line_colour': Color(rgb=(0, 0, 0)),
        'fill_colour': Color(rgb=(0, 0, 0)),
        'font_colour': Color(rgb=(0, 0, 0)),
        'text_layout': Color(rgb=(0, 0, 0)),
        'corner_radius': Cm(0.1)
    }
}

layout_attributes_test_records = {
    'layout_01': {
        'swimlane_name': 'yyy',
        'track_number': 1,
        'number_of_tracks_to_span': 1,
        'text_layout': 'shape'
    }
}

num_days_in_date_range = 365

swimlane_test_data_records = {
    'swim_01': {
        'Swimlane 1': {
            'start_track': 1,
            'end_track': 5
        }
    }
}

activity_test_data = {
    'act_01': {
        'activity_id': 12345,
        'description': 'Test activity 01',
        'activity_type': 'activity',
        'start_date': parse_date('2021-01-01'),
        'end_date': parse_date('2021-01-31'),
        'display_shape': VisualElementShape.ROUNDED_RECTANGLE,
    }
}


test_records = [
    (
        activity_test_data['act_01'],
        layout_attributes_test_records['layout_01'],
        display_attribute_test_data['disp_01'],
        display_attribute_test_data['disp_01'],
        plan_visual_config_test_data['vis_cfg_01'],
        swimlane_test_data_records['swim_01'],
        {  # Expected results record for this combination of test data
            'is_past': True,
            'is_future': False,
            'is_current': False,
            'plot_start': 0,
        }
    )
]


@ddt
class TestPlanActivity(TestCase):
    @data(*test_records)
    @unpack
    def test_create_plan_activity(
            self,
            activity_data,
            layout_data,
            display_data,
            done_display_data,
            vis_cfg,
            swimlane_data,
            exp_res
    ):
        layout_properties = ActivityLayoutAttributes(
            layout_data['swimlane_name'],
            layout_data['track_number'],
            layout_data['number_of_tracks_to_span'],
            layout_data['text_layout']
        )

        display_properties = ShapeFormatting(
            line_colour=display_data['line_colour'],
            fill_colour=display_data['line_colour'],
            corner_radius=display_data['corner_radius']
        )

        done_display_properties = ShapeFormatting(
            line_colour=done_display_data['line_colour'],
            fill_colour=done_display_data['line_colour'],
        )
        plot_driver = PlotDriver(vis_cfg)
        plot_driver.num_days_in_date_range = num_days_in_date_range

        activity = PlanActivity(
            activity_data['activity_id'],
            activity_data['description'],
            activity_data['activity_type'],
            activity_data['start_date'],
            activity_data['end_date'],
            layout_properties,
            activity_data['display_shape'],
            plot_driver,
            display_properties,
            display_properties,
            today_override=None,
            swimlane_start_track=1,
        )

        pres = Presentation()
        slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(slide_layout)
        shapes = slide.shapes

        plotted_shapes = activity.plot_ppt_shapes(shapes)

        self.assertEqual(1, len(plotted_shapes))
        plotted_shape = plotted_shapes[0]

        self.assertEqual(exp_res['is_past'], activity.is_past())
        self.assertEqual(exp_res['is_future'], activity.is_future())
        self.assertEqual(exp_res['is_current'], activity.is_current())
        self.assertEqual(exp_res['plot_start'], plotted_shape.left)
