from unittest import TestCase
from colour import Color
from ddt import ddt, unpack, data
from pptx import Presentation
from pptx.util import Cm
from source.plot_driver import PlotDriver
from source.refactor_temp.activity_layout_attributes import ActivityLayoutAttributes
from source.refactor_temp.plan_activity import PlanActivity
from source.refactor_temp.shape_formatting import ShapeFormatting
from source.refactor_temp.text_formatting import TextFormatting
from source.refactor_temp.visual_element_shape import VisualElementShape
from source.tests.testing_utilities import parse_date
from unittest.mock import Mock

visual_parameters_01 = {
    'visual_start_date': parse_date('2021-01-01'),
    'visual_end_date': parse_date('2021-07-31'),
    'visual_left': Cm(0),
    'visual_right': Cm(33.87),
}

test_case_01 = {
    'parameters': {
        'activity_start_date': parse_date('2021-01-15'),
        'activity_end_date': parse_date('2021-07-27'),
        'today': parse_date('2021-06-29'),
        'num_days_in_date_range': 212
    },
    'expected_results': {  # Calculations are from Excel
        'num_shapes': 2,
        'left_1': 805211,
        'width_1': 9489991,
        'left_2': 10295202,
        'width_2': 1667938,
    }
}


def test_case_gen():
    for expected_result in test_case_01['expected_results']:
        yield visual_parameters_01, \
              test_case_01['parameters'], \
              expected_result, \
              test_case_01['expected_results'][expected_result]


@ddt
class TestHorizontalLayout(TestCase):
    @data(*test_case_gen())
    @unpack
    def test_horizontal_layout(self, visual_parameters, activity_parameters, expected_case, expected_value):
        layout_attributes = ActivityLayoutAttributes(
            'Swimlane-01',
            1,
            1,
            "shape"
        )
        display_shape = VisualElementShape.RECTANGLE
        parameters = test_case_01['parameters']
        expected_results = test_case_01['expected_results']
        plan_visual_config = PlotDriver(
            {
                'top': Cm(3.86),
                'left': visual_parameters['visual_left'],
                'bottom': Cm(20),
                'right': Cm(33.87),
                'track_height': Cm(0.6),
                'track_gap': Cm(0.2),
                'activity_text_width': Cm(5),
                'milestone_text_width': Cm(5),
                'text_margin': Cm(0.1),
                'milestone_width': Cm(0.4),
                'activity_shape': 'rectangle',
                'milestone_shape': 'diamond',
                'min_start_date': parse_date('2021-01-01'),
                'max_end_date': parse_date('2021-07-31')
            }
        )
        plan_visual_config.num_days_in_date_range = parameters['num_days_in_date_range']
        text_formatting = TextFormatting()
        formatting_1 = ShapeFormatting(
            Color(rgb=(1,1,1)),
            Color(rgb=(1,1,1)),
            None,
            text_formatting,
        )

        activity = PlanActivity(
            12345,
            'Dummy',
            'bar',
            parameters['activity_start_date'],
            parameters['activity_end_date'],
            layout_attributes,
            display_shape,
            plan_visual_config,
            formatting_1,
            formatting_1,
            parameters['today'],
            1
        )

        pres = Presentation()
        slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(slide_layout)
        shapes = slide.shapes

        shapes = activity.plot_ppt_shapes(shapes)
        actual_left_1 = shapes[0].left
        actual_width_1 = shapes[0].width
        actual_left_2 = shapes[1].left
        actual_width_2 = shapes[1].width

        if expected_case == 'num_shapes':
            self.assertEqual(expected_results['num_shapes'], len(shapes))
        elif expected_case == 'left_1':
            self.assertEqual(expected_results['left_1'], actual_left_1)
        elif expected_case == 'width_1':
            self.assertEqual(expected_results['width_1'], actual_width_1)
        elif expected_case == 'left_2':
            self.assertEqual(expected_results['left_2'], actual_left_2)
        elif expected_case == 'width_2':
            self.assertEqual(expected_results['width_2'], actual_width_2)