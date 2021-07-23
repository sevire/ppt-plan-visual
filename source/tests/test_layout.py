from unittest import TestCase
from colour import Color
from ddt import ddt, unpack, data
from pptx import Presentation
from pptx.util import Cm
from source.visualiser.plot_driver import PlotDriver
from source.visualiser.activity_layout_attributes import ActivityLayoutAttributes
from source.visualiser.plan_activity import PlanActivity
from source.visualiser.shape_formatting import ShapeFormatting
from source.visualiser.text_formatting import TextFormatting
from source.visualiser.visual_element_shape import VisualElementShape
from source.tests.testing_utilities import parse_date

visual_parameters_01 = {
    'visual_start_date': parse_date('2021-01-01'),
    'visual_end_date': parse_date('2021-07-31'),
    'visual_left': Cm(0),
    'visual_right': Cm(33.87),
}

test_cases = [
    {
        'parameters': {
            'activity_start_date': parse_date('2021-08-11'),
            'activity_end_date': parse_date('2021-12-14'),
            'today': parse_date('2021-07-01'),
            'num_days_in_date_range': 212
        },
        'expected_results': {  # Calculations are from Excel
            'num_shapes': 1,
            'left': 12768351,
            'width': 7246902,
            'top': 1965600
        }
    },
    {
        'parameters': {
            'activity_start_date': parse_date('2021-01-15'),
            'activity_end_date': parse_date('2021-07-27'),
            'today': parse_date('2021-06-29'),
            'num_days_in_date_range': 212
        },
        'expected_results': {  # Calculations are from Excel
            'num_shapes': 2,
            'left_1': 805211,
            'width_1': 9489991,
            'left_2': 10295202,
            'width_2': 1667938,
        }
    }
]


def test_case_gen():
    for test_case in test_cases:
        for expected_result in test_case['expected_results']:
            yield visual_parameters_01, \
                  test_case['parameters'], \
                  expected_result, \
                  test_case['expected_results'][expected_result]


@ddt
class TestLayout(TestCase):
    @data(*test_case_gen())
    @unpack
    def test_horizontal_layout(self, visual_parameters, activity_parameters, expected_case, expected_value):
        layout_attributes = ActivityLayoutAttributes(
            'Swimlane-01',
            1,
            1,
            "shape"
        )
        display_shape = VisualElementShape.RECTANGLE
        # parameters = test_case_01['parameters']
        # expected_results = test_case_01['expected_results']
        plan_visual_config = PlotDriver(
            {
                'top': Cm(3.86),
                'left': visual_parameters['visual_left'],
                'bottom': Cm(20),
                'right': Cm(33.87),
                'track_height': Cm(0.6),
                'track_gap': Cm(0.2),
                'activity_text_width': Cm(5),
                'milestone_text_width': Cm(5),
                'text_margin': Cm(0.1),
                'milestone_width': Cm(0.4),
                'activity_shape': 'rectangle',
                'milestone_shape': 'diamond',
                'min_start_date': parse_date('2021-01-01'),
                'max_end_date': parse_date('2021-07-31')
            }
        )
        plan_visual_config.num_days_in_date_range = activity_parameters['num_days_in_date_range']
        text_formatting = TextFormatting()
        formatting_1 = ShapeFormatting(
            Color(rgb=(1,1,1)),
            Color(rgb=(1,1,1)),
            None,
            text_formatting,
        )

        activity = PlanActivity(
            12345,
            'Dummy',
            'bar',
            activity_parameters['activity_start_date'],
            activity_parameters['activity_end_date'],
            layout_attributes,
            display_shape,
            plan_visual_config,
            formatting_1,
            formatting_1,
            activity_parameters['today'],
            1
        )

        pres = Presentation()
        slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(slide_layout)
        shapes = slide.shapes

        shapes = activity.plot_ppt_shapes(shapes)

        if expected_case == 'num_shapes':
            self.assertEqual(expected_value, len(shapes))
        elif expected_case == 'left_1':
            self.assertEqual(expected_value, shapes[0].left)
        elif expected_case == 'width_1':
            self.assertEqual(expected_value, shapes[0].width)
        elif expected_case == 'left_2':
            self.assertEqual(expected_value, shapes[1].left)
        elif expected_case == 'width_2':
            self.assertEqual(expected_value, shapes[1].width)