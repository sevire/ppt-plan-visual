from unittest import TestCase

from ddt import ddt, data, unpack
from pptx import Presentation

from source.tests.test_resources.unit_test_01.expected_results import input_files_01, expected_results_01, today
from source.visualiser.plan_visualiser import PlanVisualiser


def plan_test_case_generator():
    """
    Generator to break up tests so that we can have one assert per test.

    Each activity will generate 2 or three shapes, one or two graphic shapes and one text shape.  All the shapes
    have positioning and formatting which needs to be checked.

    We need to generate expected results individually under the following structure:
    - 2 or 3 shapes plotted, depending upon whether the shape was plotted with a done/to-do split or not.
      - Fields for 1 or 2 objects representing the activity bar or milestone graphic.
      - Fields for the text label for the activity.

    :return:
    """
    field_names = ['top', 'left', 'width']  # Drives which field is being tested in given test
    for activity_seq_num, activity_exp_result in enumerate(expected_results_01["plan_data"]):
        activity_type, activity_text, shape_data = activity_exp_result
        num_shapes = len(shape_data)  # Will be 2 or 3
        yield activity_seq_num, None, 'num_shapes', num_shapes

        # Yield shape 1 field expected results
        result_pairs = zip(field_names, shape_data[0])
        for result in result_pairs:
            # This will feed key fields to the tewt to be able to generate the right actual result to test against.
            yield activity_seq_num, 'graphic_shape_1', result[0], result[1]

        # Yield shape 2 field expected results if there are two graphic shapes for this activity
        if num_shapes == 3:
            result_pairs = zip(field_names, shape_data[1])
            for result in result_pairs:
                # This will feed key fields to the tewt to be able to generate the right actual result to test against.
                yield activity_seq_num, 'graphic_shape_2', result[0], result[1]

        # Yield text field expected results
        # First yield the text field (only for the text shape, not for graphic shapes)
        yield activity_seq_num, 'text_shape', 'text', activity_text
        text_exp_res_index = num_shapes - 1  # Text field data will be the 2nd or 3rd entry depending upon num_shapes
        result_pairs = zip(field_names, shape_data[text_exp_res_index])
        for result in result_pairs:
            # This will feed key fields to the tewt to be able to generate the right actual result to test against.
            yield activity_seq_num, 'text_shape', result[0], result[1]


@ddt
class TestComprehensive01(TestCase):
    def setUp(self) -> None:
        """
        Create plan object which drives tests.

        :return:
        """
        excel_plan_file = input_files_01["excel_plan_file"]
        excel_config_workbook = input_files_01["visual_config"]
        ppt_template_file = input_files_01["ppt_template"]
        excel_plan_sheet = input_files_01["plan_sheet_name"]

        self.visualiser = PlanVisualiser.from_excel(
            excel_plan_file,
            excel_config_workbook,
            ppt_template_file,
            excel_plan_sheet)

        # Need to plot swimlanes as that calculates where swimlanes are positioned vertically which then impacts upon
        # positioning of plan elements.
        self.visualiser.plot_swimlanes(self.visualiser.format_config)

        # We need a shapes object to call plot shapes with (even though we aren't actually plotting to PPT!)
        SLD_LAYOUT_TITLE_AND_CONTENT = 1
        prs = Presentation()
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = prs.slides.add_slide(slide_layout)

        self.shapes = slide.shapes

    def compare_shape_data(self, shape, top, left):
        self.assertEqual(top, shape.top)
        self.assertEqual(left, shape.left)

    @data(*plan_test_case_generator())
    @unpack
    def test_plan_01(self, activity_num, shape_to_test, field_name, expected_value):
        """
        Iterate through elements of plan in entry order which corresponds with the order of expected results.
        Call plot_ppt_shapes and then check the shape objects returned against expected results.

        :return:
        """
        expected_plan_results = expected_results_01["plan_data"]
        self.assertEqual(len(expected_plan_results),len(self.visualiser.plan_data))

        activity = self.visualiser.plan_data[activity_num]
        activity.today_override = today
        activity.swimlane_start_track = self.visualiser.swimlane_data[activity.activity_layout_attributes.swimlane_name]['start_track']
        shapes = activity.plot_ppt_shapes(self.visualiser.shapes)
        num_shapes = len(shapes)
        shape = None
        if shape_to_test is None:
            # This test cases is not at shape level (probably number of shapes)
            shape = None
        if shape_to_test == 'graphic_shape_1':
            shape = shapes[0]
        elif shape_to_test == 'graphic_shape_2':
            shape = shapes[1]
        elif shape_to_test == 'text_shape':
            shape = shapes [-1]  # Text will always be the last shape

        if field_name == 'num_shapes':
            self.assertEqual(expected_value, num_shapes)
        elif field_name == 'top':
            self.assertEqual(expected_value, shape.top)
        elif field_name == 'left':
            self.assertEqual(expected_value, shape.left)
        elif field_name == 'text':
            self.assertEqual(expected_value, shape.text)
        elif field_name == 'width':
            self.assertEqual(expected_value, shape.width)
        else:
            self.fail("Unknown test parameters")
