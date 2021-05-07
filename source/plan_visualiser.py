import logging
import os
from calendar import month_name
from datetime import date
from functools import reduce

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from source.excel_plan import ExcelPlan
from source.plot_driver import PlotDriver
from source.utilities import get_path_name_ext, SwimlaneManager, first_day_of_month, iterate_months, \
    num_months_between_dates, last_day_of_month, is_current, is_nan, is_future, is_past, day_increment

root_logger = logging.getLogger()


class PlanVisualiser:
    """
    Object which manages the creation of a PPT slide with a visual representation
    of a project plan (or similar).

    The supplied data

    :param plan_data:
    """

    def __init__(self, plan_data, plot_config, format_config, template_path, swimlanes):
        # The actual plan data with activities and milestones, start/finish dates etc.
        self.plan_data = plan_data

        # Data to define plot area for elements, tracks etc.
        self.plot_config = plot_config

        # Data with pre-determined formatting properties to apply to elements.
        self.format_config = format_config

        # # Config to drive slide level objects such as the swimlane rectangle shapes as background.
        # self.slide_level_config = slide_level_config
        #
        self.template = template_path
        folder, base, ext = get_path_name_ext(template_path)
        self.slides_out_path = os.path.join(folder, base + '_out' + ext)

        self.prs = Presentation(template_path)

        visual_slide = self.prs.slides[0]  # Assume there is one slide and that's where we will place the visual

        self.shapes = visual_slide.shapes
        self.plot_driver = PlotDriver(plot_config)

        self.align_months()

        # We can't be sure of knowing the number of days in the range until align_months has been called, so do the
        # calculation here, not in PlotDriver (where it used to be)
        self.plot_driver.num_days_in_date_range = self.plot_driver.max_end_date.toordinal() - self.plot_driver.min_start_date.toordinal() + 1

        self.swimlanes = swimlanes
        self.swimlane_data = self.extract_swimlane_data()

    def plot_slide(self):
        """
        Opens a supplied template file in order to allow consistency with other slides in a deck.

        Then creates a new slide with a plotted plan visual based on the data read in to the object.

        Then writes the one-slide deck to a different filename in the same folder.

        :return:
        """

        self.plot_swimlanes(self.format_config)
        self.plot_month_bar()

        root_logger.info(f'Plotting {len(self.plan_data)} elements')

        for plotable_element in self.plan_data:
            start = plotable_element['start_date']
            end = plotable_element['end_date']
            description = plotable_element['description']

            root_logger.debug(f'Plotting activity: [{description:40.40}], start: {start}, end: {end}')

            swimlane = plotable_element['swimlane']
            track_num = plotable_element['track_num']
            num_tracks = plotable_element['bar_height_in_tracks']

            shape_format_name = plotable_element['format_properties']
            done_shape_format_name = plotable_element['done_format_properties']
            format_data = self.format_config[shape_format_name]
            done_format_data = None if is_nan(done_shape_format_name) else self.format_config[done_shape_format_name]

            text_layout = plotable_element['text_layout']
            if plotable_element['type'] == 'bar':
                self.plot_activity(description, start, end, swimlane, track_num, num_tracks, format_data, text_layout, done_format_data)

            elif plotable_element['type'] == 'milestone':
                self.plot_milestone(description, start, swimlane, track_num, format_data, text_layout)

        self.plot_vertical_line(date.today())
        self.prs.save(self.slides_out_path)

    def plot_activity(self, activity_description, start_date, end_date, swimlane, track_number, num_tracks,
                      todo_properties, text_layout, done_properties):
        """

        :param activity_description:
        :param start_date:
        :param end_date:
        :param swimlane:
        :param track_number:
        :param num_tracks:
        :param todo_properties:
        :param text_layout:
        :param done_properties:
        :return:

        Plots an activity on the slide (rather than a milestone)

        Plots rectangle for the activity of the appropriate length based on the length of the activity and
        places on the right track at the right horizontal position based on start date.

        If a done_properties value has been provided then this is a task where we need to indicate any past
        part of the activity as "Done" using the format provided to indicate that (typically this would be blue).

        So if the activity is current (today's date is between start and end) we have to plot two
        boxes of different colours, but if the activity is completely in the past we plot one box with the
        done formatting.
        """
        swimlane_start = self.swimlane_data[swimlane]['start_track']

        top = self.plot_driver.track_number_to_y_coordinate(swimlane_start + track_number - 1)
        height = self.plot_driver.height_of_track(num_tracks)

        # There are three cases to consider.
        # - There are no done_properties provided or the activity is wholly in the future
        #   * One rectangle with todo_properties
        # - There are done_properties provided and the activity is current
        #   * One rectangle for done period with done_properties
        #   * One rectangle for still to do period with todo_properties
        # - There are done_properties provided and the activity is wholly in the past
        #   * One rectangle with todo_properties

        if done_properties is None or is_future(start_date, end_date):
            left, right, width = self.shape_parameters(start_date, end_date)

            self.plot_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, todo_properties)

        elif done_properties is not None and is_current(start_date, end_date):
            # One rectangle for done period with done_properties
            # One rectangle for still to do period with todo_properties

            done_left, done_right, done_width = self.shape_parameters(start_date, date.today(), gap_flag=False)
            self.plot_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, done_left, top, done_width, height, done_properties)

            todo_left, todo_right, todo_width = self.shape_parameters(date.today(), end_date, gap_flag=True)
            self.plot_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, todo_left, top, todo_width, height, todo_properties)
        elif done_properties is not None and is_past(start_date, end_date):
            # NOTE this test should always be true if previous two tests false but this is belt and braces

            # One rectangle with todo_properties
            left, right, width = self.shape_parameters(start_date, date.today())

            self.plot_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, done_properties)

        else:
            # Shouldn't get to here so raise an exception
            raise Exception("Shouldn't get here")

        # Text is same plot whether we are plotting as two activities or one
        left, right, width = self.shape_parameters(start_date, end_date)
        self.plot_text_for_shape(left, top, width, height, activity_description, todo_properties, text_layout)

    def plot_milestone(self, milestone_description, start_date, swimlane, track_number, properties, text_layout):
        swimlane_start = self.swimlane_data[swimlane]['start_track']
        milestone_width = self.plot_config['milestone_width']
        milestone_height = self.plot_config['track_height']

        left = self.plot_driver.date_to_x_coordinate(start_date) - milestone_width / 2
        top = self.plot_driver.track_number_to_y_coordinate(swimlane_start + track_number - 1)

        shape = self.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.DIAMOND, left, top, milestone_width, milestone_height
        )

        self.shape_fill(shape, properties)
        self.shape_line(shape, properties)

        if text_layout == "Right":
            milestone_text_width = self.plot_config['milestone_text_width']
            milestone_text_left = left + milestone_width
            properties['text_align'] = 'left'
            self.plot_text(milestone_description, milestone_text_left, top, milestone_text_width, milestone_height, properties)
        else:  # Default is left
            milestone_text_width = self.plot_config['milestone_text_width']
            milestone_text_left = left - milestone_text_width
            properties['text_align'] = 'right'
            self.plot_text(milestone_description, milestone_text_left, top, milestone_text_width, milestone_height, properties)

    def plot_shape(self, shape_type, left, top, width, height, shape_properties):
        shape = self.shapes.add_shape(
            shape_type, left, top, width, height
        )

        # Adjust rounded corner radius
        if shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
            target_radius = shape_properties['corner_radius']
            adjustment_value = target_radius / height
            shape.adjustments[0] = adjustment_value

        self.shape_fill(shape, shape_properties)
        self.shape_line(shape, shape_properties)

    def plot_text_for_shape(self, left, top, width, height, text, shape_properties, text_layout):
        activity_text_width = self.plot_config['activity_text_width']
        if text_layout == 'Left':
            # Extend the text to the left so that if overflows to the left of the shape.
            adjust_width = max(width, activity_text_width)
            text_left = left + width - adjust_width
            text_width = activity_text_width
            shape_properties['text_align'] = 'right'
            self.plot_text(text, text_left, top, text_width, height, shape_properties)
        elif text_layout == 'Right':
            # Extend text to the right so that it overflows to the right of the shape
            adjust_width = max(width, activity_text_width)
            text_left = left
            text_width = adjust_width
            shape_properties['text_align'] = 'left'
            self.plot_text(text, text_left, top, text_width, height, shape_properties)
        else:  # Apply default which is "Shape"
            # Standard positioning, text will align exactly with the shape
            shape_properties['text_align'] = 'centre'
            self.plot_text(text, left, top, width, height, shape_properties)

    def plot_text(self, text, left, top, width, height, format_data):

        shape = self.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
        )
        shape.fill.background()
        shape.line.fill.background()

        self.add_text_to_shape(shape, text, format_data)

    def add_text_to_shape(self, shape, text, format_data):
        text_frame = shape.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        vertical = format_data['text_vertical_align']
        if vertical == "top":
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        elif vertical == "bottom":
            text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text

        self.text_format(paragraph, run, format_data)

    @staticmethod
    def shape_fill(shape, format_data):
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*format_data['fill_rgb'])

    def text_format(self, paragraph, run, format_data):
        """
        Hard code text formatting (for now) as is likely to be constant.

        :param paragraph:
        :param run:
        :return:
        """
        font = run.font
        paragraph.line_spacing = 0.8
        paragraph.alignment = self.text_alignment(format_data['text_align'])

        font.name = 'Calibri'
        font.size = format_data['font_size']
        font.bold = format_data['font_bold']
        font.italic = format_data['font_italic']
        font.color.rgb = RGBColor(*format_data['font_colour_rgb'])

    @staticmethod
    def shape_line(shape, format_data):
        line = shape.line
        line.color.rgb = RGBColor(*format_data['line_rgb'])

    @staticmethod
    def text_alignment(format_text_align):
        """
        Takes the alignment field from the element formatting data and converts to the appropriate value for the
        pptx-python setting.

        :param format_text_align:
        :return:
        """
        if format_text_align == "left":
            return PP_ALIGN.LEFT
        elif format_text_align == "right":
            return PP_ALIGN.RIGHT
        elif format_text_align == "centre":
            return PP_ALIGN.CENTER
        else:
            # Default to centre
            return PP_ALIGN.CENTER

    def plot_swimlanes(self, format_data):
        """
        plot a background rectangle for each swimline with the width of the plot area, the height of the swimlane, and
        positioned to coincide with the tracks within the swimlane.

        Also alternate colours for each swimlane

        :return:
        """

        for row, swimlane in enumerate(self.swimlane_data):
            row_number = row + 1

            start_track = self.swimlane_data[swimlane]['start_track']
            end_track = self.swimlane_data[swimlane]['end_track']
            top = self.plot_driver.track_number_to_y_coordinate(start_track)

            # Need to adjust to start between track end and track start, unless this is the first row
            if row_number > 1:
                top -= (self.plot_config['track_gap'] / 2)
            bottom = self.plot_driver.track_number_to_y_coordinate(end_track) + self.plot_config['track_height'] + (self.plot_config['track_gap'] / 2)
            left = self.plot_config['left']
            width = self.plot_config['right'] - self.plot_config['left']
            height = bottom - top

            shape = self.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )

            # For the purposes of this decision, the first row is 1 (odd)
            if row_number % 2 == 0:
                format_info = format_data['swimlane_format_even']
            else:
                format_info = format_data['swimlane_format_odd']
            # Hard code alignment for now.  May need to re-visit
            format_info['text_align'] = 'left'
            # ToDo: Add configuration of horizontal alignment for swimlanes

            self.shape_fill(shape, format_info)
            self.shape_line(shape, format_info)
            self.add_text_to_shape(shape, swimlane, format_info)

    def plot_month_bar(self):
        """
        Create a rectangle for each month on the timeline, driven by the configured start and end date for the plan.
        The rectangles will have a configured height and each will be the width corresponding to the number of days
        in the month.

        :return:
        """
        first_of_start_month = first_day_of_month(self.plot_driver.min_start_date)
        first_of_end_month = first_day_of_month(self.plot_driver.max_end_date)

        for month_index, month_start_date in enumerate(iterate_months(first_of_start_month, num_months_between_dates(first_of_start_month, first_of_end_month))):
            # left = self.plot_driver.date_to_x_coordinate(month_start_date)
            # right = self.plot_driver.date_to_x_coordinate(day_increment(last_day_of_month(month_start_date), 1))
            # width = right - left

            left, right, width = self.shape_parameters(month_start_date, last_day_of_month(month_start_date), gap_flag=False)

            height = (self.plot_config['track_height'] * 1)
            top = self.plot_config['top'] - height

            if month_index % 2 == 1:
                shape_format = self.format_config['month_shape_format_odd']
            else:
                shape_format = self.format_config['month_shape_format_even']

            month = month_name[month_start_date.month][:3]

            self.plot_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height, shape_format)
            self.plot_text_for_shape(left, top, width, height, month, shape_format, 'shape')

    @classmethod
    def from_excel_plan(cls, plan_data_excel_path, plan_data_sheet_name):
        extracted_plan_data = ExcelPlan.read_plan_data(plan_data_excel_path, plan_data_sheet_name)

        return extracted_plan_data

    def extract_swimlane_data(self):
        """
        After some thought am taking a very simple approach here.

        Just go through each activity on the plan and:
        - Check which swimlane it is for
        - Look at the track number and height in tracks and therefore work out the bottom track number
        - If the bottom track number is larger than that already recorded, update to reflect.
        - Once you have checked every activity, the entry in each swimlane will have recorded the number of tracks
          required for that swimlane.
        - We can then go back and calculate the start and end track for each swimlane which is what the plot method
          needs.  If a specific swimlane order is dictated then that is used here.  Where a swimlane doesn't appear in
          the ordering then we just add it to the end.

          So we will end up with a dict, with one entry for each (named) swimlane, and against each swimlane we will
          see the start track number and the end track number.

        :return:
        """

        swimlane_data = {}
        swimlane_manager = SwimlaneManager(self.swimlanes)
        for record in self.plan_data:
            swimlane = record['swimlane']
            track_num_high = record['track_num'] + record['bar_height_in_tracks'] - 1
            if swimlane not in swimlane_data:
                # Adding record for this swimlane. Also remember ordering as is important later.
                swimlane_record = {
                    'swimlane_number': swimlane_manager.get_swimlane_number(swimlane),
                    'highest_track_within_lane': track_num_high
                }
                swimlane_data[swimlane] = swimlane_record
            else:
                swimlane_record = swimlane_data[swimlane]
                if track_num_high > swimlane_record['highest_track_within_lane']:
                    swimlane_record['highest_track_within_lane'] = track_num_high

        # We now have a dict containing a record for each swimlane of lowest and highest relative track number used.
        # Can now calculate the start and end track number for each swimlane - in order that lanes were encountered.

        swimlane_plot_data = {}
        end_track = 0

        # Create list of swimlane data entries in order of swimlane number
        sorted_entries = sorted(swimlane_data.items(), key=lambda x: x[1]['swimlane_number'])

        # Calculate start and end track number for each swimlane based on number of tracks required for swimlane
        for swimlane, swimlane_entry in sorted_entries:
            start_track = end_track + 1
            end_track = start_track + swimlane_entry['highest_track_within_lane'] - 1
            swimlane_plot_data[swimlane] = {
                'start_track': start_track,
                'end_track': end_track
            }
        return swimlane_plot_data

    def align_months(self):
        """
        If earliest or latest dates haven't been specified explicitly, then calculate from plan data.
        Then adjust to be first and last days of month respectively to ensure that month bar
        aligns with configured plot area.

        :return:
        """

        if pd.isnull(self.plot_driver.min_start_date):

            start_dates = [record['start_date'] for record in self.plan_data]
            self.plot_driver.min_start_date = reduce(lambda min_date, start_date: start_date if start_date < min_date else min_date, start_dates)

        if pd.isnull(self.plot_driver.max_end_date):
            end_dates = [record['end_date'] for record in self.plan_data]
            self.plot_driver.max_end_date = reduce(lambda max_date, end_date: end_date if not pd.isnull(end_date) and end_date > max_date else max_date, end_dates)

        # Regardless of whether start and end dates have been configured, we need to align with whole month
        self.plot_driver.min_start_date = first_day_of_month(self.plot_driver.min_start_date)
        self.plot_driver.max_end_date = last_day_of_month(self.plot_driver.max_end_date)

    def shape_parameters(self, start_date, end_date, gap_flag=True):
        """
        Calculates the (x-axis) parameters required for plotting a shape based on a start and end date.
        To allow a visible but very small gap between activities, the right hand edge is brought in by a fixed number of
        units, but to allow for cases where this isn't desirable (e.g. when plotting month bars or the done and to do
        parts of an activity, when gap_flag is False this adjustment won't be made.

        :param start_date:
        :param end_date:
        :param gap_flag:
        :return:
        """
        if gap_flag is True:
            gap_increment = 15000
        else:
            gap_increment = 0

        left = self.plot_driver.date_to_x_coordinate(start_date)
        right = self.plot_driver.date_to_x_coordinate(day_increment(end_date, 1)) - gap_increment
        width = right - left

        return left, right, width

    def plot_vertical_line(self, date):
        x = self.plot_driver.date_to_x_coordinate(date)
        top = self.plot_config['top']
        bottom = self.plot_config['bottom']
        line = self.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x, top, x, bottom)
        today_line_format = self.format_config['today_line']
        today_line_colour = today_line_format['line_rgb']
        line.line.color.rgb = RGBColor(*today_line_colour)
