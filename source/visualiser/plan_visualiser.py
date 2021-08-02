import logging
import os
from calendar import month_name
from datetime import date
from functools import reduce
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR

from source.visualiser.excel_config import ExcelPlotConfig, ExcelFormatConfig, ExcelSwimlaneConfig
from source.visualiser.excel_plan import ExcelPlan
from source.visualiser.plan_activity import PlanActivity
from source.visualiser.plot_driver import PlotDriver
from source.visualiser.plotable_element import PlotableElement
from source.visualiser.shape_formatting import ShapeFormatting
from source.visualiser.text_formatting import TextFormatting
from source.visualiser.visual_element_shape import VisualElementShape
from source.visualiser.utilities import get_path_name_ext, SwimlaneManager, first_day_of_month, iterate_months, \
    num_months_between_dates, last_day_of_month

root_logger = logging.getLogger()


class PlanVisualiser:
    """
    Object which manages the creation of a PPT slide with a visual representation
    of a project plan (or similar).

    The supplied data

    :param plan_data:
    """

    def __init__(
            self,
            plan_data: List[PlanActivity],
            plot_config: PlotDriver,
            format_config: dict,
            template_path: str,
            swimlanes: List[dict]):
        # The actual plan data with activities and milestones, start/finish dates etc.
        self.plan_data = plan_data

        # Data to define plot area for elements, tracks etc.
        self.plot_config = plot_config

        # Data with pre-determined formatting properties to apply to elements.
        self.format_config = format_config

        # # Config to drive slide level objects such as the swimlane rectangle shapes as background.
        # self.slide_level_config = slide_level_config
        #
        self.template = template_path
        folder, base, ext = get_path_name_ext(template_path)
        self.slides_out_path = os.path.join(folder, base + '_out' + ext)

        self.prs = Presentation(template_path)

        visual_slide = self.prs.slides[0]  # Assume there is one slide and that's where we will place the visual

        self.shapes = visual_slide.shapes
        self.plot_driver = plot_config

        self.align_months()

        # We can't be sure of knowing the number of days in the range until align_months has been called, so do the
        # calculation here, not in PlotDriver (where it used to be)
        self.plot_driver.num_days_in_date_range = self.plot_driver.max_end_date.toordinal() - self.plot_driver.min_start_date.toordinal() + 1

        self.swimlanes = swimlanes
        self.swimlane_data = self.extract_swimlane_data()

    @classmethod
    def from_excel(cls, excel_plan_file, excel_config_workbook, ppt_template_file, excel_plan_sheet=None):
        """
        Reads plan and configuration information from Excel workbooks and then creates instance of PlanVisualiser

        :return:
        """
        print("Plan Visualiser - starting...")
        print("Initiating logging")
        root_logger.debug('Plan to PowerPoint plotting programme starting...')
        root_logger.info(f"Running from IDE, using fixed arguments")

        root_logger.info(f'Using plan data from {excel_plan_file}')

        plot_config_object = ExcelPlotConfig(excel_config_workbook, excel_sheet='PlotConfig')
        plot_area_config = plot_config_object.parse_plot_config()

        excel_format_config_object = ExcelFormatConfig(excel_config_workbook, excel_sheet='FormatConfig')
        shape_config = excel_format_config_object.parse_format_config()

        plan_data = ExcelPlan.read_plan_data(
            excel_plan_file,
            excel_plan_sheet,
            shape_config,
            plot_area_config
        )

        swimlane_config_object = ExcelSwimlaneConfig(excel_config_workbook, excel_sheet='Swimlanes')
        swimlanes = swimlane_config_object.parse_swimlane_config()

        return cls(plan_data, plot_area_config, shape_config, ppt_template_file, swimlanes)

    def plot_slide(self):
        """
        Opens a supplied template file in order to allow consistency with other slides in a deck.

        Then creates a new slide with a plotted plan visual based on the data read in to the object.

        Then writes the one-slide deck to a different filename in the same folder.

        :return:
        """

        self.plot_swimlanes(self.format_config)
        self.plot_month_bar()

        root_logger.info(f'Plotting {len(self.plan_data)} elements')

        for activity in self.plan_data:
            start = activity.start_date
            end = activity.end_date
            description = activity.description

            root_logger.debug(f'Plotting activity: [{description:40.40}], start: {start}, end: {end}')

            activity.swimlane_start_track = self.swimlane_data[activity.activity_layout_attributes.swimlane_name]['start_track']
            activity.plot_ppt_shapes(self.shapes)

        self.plot_vertical_line(date.today())
        self.prs.save(self.slides_out_path)

    def plot_text_for_shape(self, left, top, width, height, text, shape_properties, text_layout):
        activity_text_width = self.plot_config.min_activity_text_width
        if text_layout == 'Left':
            # Extend the text to the left so that if overflows to the left of the shape.
            adjust_width = max(width, activity_text_width)
            text_left = left + width - adjust_width
            text_width = activity_text_width
            shape_properties.text_align = 'right'
            self.plot_text(text, text_left, top, text_width, height, shape_properties, text_layout)
        elif text_layout == 'Right':
            # Extend text to the right so that it overflows to the right of the shape
            adjust_width = max(width, activity_text_width)
            text_left = left
            text_width = adjust_width
            shape_properties['text_align'] = 'left'
            self.plot_text(text, text_left, top, text_width, height, shape_properties, text_layout)
        else:  # Apply default which is "Shape"
            # Standard positioning, text will align exactly with the shape
            shape_properties['text_align'] = 'centre'
            self.plot_text(text, left, top, width, height, shape_properties, text_layout)

    def plot_text(self, text, left, top, width, height, format_data, text_layout):

        shape = self.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
        )
        shape.fill.background()
        shape.line.fill.background()

        self.add_text_to_shape(shape, text, format_data, text_layout)

    def add_text_to_shape(self, shape, text, format_data, text_layout):
        text_frame = shape.text_frame

        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0

        # Adjust text margin depending upon positioning. To help readability by having small gap
        if text_layout == "swimlane":
            # Hard-coded case for swimlanes as isn't driven by configuration
            text_frame.margin_top = self.plot_config.text_margin
            text_frame.margin_left = self.plot_config.text_margin
        if text_layout == "Left":
            text_frame.margin_right = self.plot_config.text_margin
        elif text_layout == "Right":
            text_frame.margin_left = self.plot_config.text_margin

        vertical = format_data.text_vertical_align
        if vertical == "top":
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        elif vertical == "bottom":
            text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text

        self.text_format(paragraph, run, format_data)

    @staticmethod
    def shape_fill(shape, format_data):
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*format_data['fill_rgb'])

    def text_format(self, paragraph, run, format_data):
        """
        Hard code text formatting (for now) as is likely to be constant.

        :param paragraph:
        :param run:
        :param format_data:
        :return:
        """
        font = run.font
        paragraph.line_spacing = 0.8
        paragraph.alignment = self.text_alignment(format_data['text_align'])

        font.name = 'Calibri'
        font.size = format_data['font_size']
        font.bold = format_data['font_bold']
        font.italic = format_data['font_italic']
        font.color.rgb = RGBColor(*format_data['font_colour_rgb'])

    @staticmethod
    def shape_line(shape, format_data):
        line = shape.line
        line.color.rgb = RGBColor(*format_data['line_rgb'])

    @staticmethod
    def text_alignment(format_text_align):
        """
        Takes the alignment field from the element formatting data and converts to the appropriate value for the
        pptx-python setting.

        :param format_text_align:
        :return:
        """
        if format_text_align == "left":
            return PP_ALIGN.LEFT
        elif format_text_align == "right":
            return PP_ALIGN.RIGHT
        elif format_text_align == "centre":
            return PP_ALIGN.CENTER
        else:
            # Default to centre
            return PP_ALIGN.CENTER

    def plot_swimlanes(self, format_data):
        """
        plot a background rectangle for each swimline with the width of the plot area, the height of the swimlane, and
        positioned to coincide with the tracks within the swimlane.

        Also alternate colours for each swimlane

        :return:
        """

        for row, swimlane in enumerate(self.swimlane_data):
            row_number = row + 1

            start_track = self.swimlane_data[swimlane]['start_track']
            end_track = self.swimlane_data[swimlane]['end_track']
            top = self.plot_driver.track_number_to_y_coordinate(start_track)

            # Need to adjust to start between track end and track start, unless this is the first row
            if row_number > 1:
                top -= (self.plot_config.track_gap / 2)
            bottom = self.plot_driver.track_number_to_y_coordinate(end_track) + self.plot_config.track_height + (
                        self.plot_config.track_gap / 2)
            left = self.plot_config.left
            right = self.plot_config.right
            width = right - self.plot_config.left
            height = bottom - top

            # For the purposes of this decision, the first row is 1 (odd)
            if row_number % 2 == 0:
                format_info = format_data['swimlane_format_even']
            else:
                format_info = format_data['swimlane_format_odd']
            # Hard code alignment for now.  May need to re-visit
            format_info['text_align'] = 'left'
            shape_formatting = ShapeFormatting.from_dict(format_info, self.plot_config)
            # ToDo: Add configuration of horizontal alignment for swimlanes

            text_formatting = TextFormatting(
                vertical_align='top',
                horizontal_align='left',
                margin_top=self.plot_config.text_margin,
                margin_bottom=self.plot_config.text_margin,
                margin_left=self.plot_config.text_margin,
                margin_right=self.plot_config.text_margin,
                font_colour=format_info['font_colour_rgb']
            )

            plottable = PlotableElement(
                VisualElementShape.RECTANGLE,
                top,
                left,
                bottom,
                right,
                shape_formatting,
                swimlane,
                text_formatting
            )
            plottable.plot_ppt(self.shapes)

    def plot_month_bar(self):
        """
        Create a rectangle for each month on the timeline, driven by the configured start and end date for the plan.
        The rectangles will have a configured height and each will be the width corresponding to the number of days
        in the month.

        :return:
        """
        first_of_start_month = first_day_of_month(self.plot_driver.min_start_date)
        first_of_end_month = first_day_of_month(self.plot_driver.max_end_date)

        for month_index, month_start_date in enumerate(iterate_months(first_of_start_month,
                                                                      num_months_between_dates(first_of_start_month,
                                                                                               first_of_end_month))):
            left, right, width = self.plot_driver.shape_parameters(month_start_date,
                                                                   last_day_of_month(month_start_date), gap_flag=False)

            height = (self.plot_config.track_height * 1)
            top = self.plot_config.top - height

            if month_index % 2 == 1:
                shape_format = self.format_config['month_shape_format_odd']
            else:
                shape_format = self.format_config['month_shape_format_even']

            shape_formatting = ShapeFormatting.from_dict(shape_format, self.plot_config)
            text_formatting = TextFormatting()

            month = month_name[month_start_date.month][:3]
            bottom = top + height
            right = left + width

            plotable = PlotableElement(
                VisualElementShape.RECTANGLE,
                top, left, bottom, right,
                shape_formatting,
                month,
                text_formatting)

            plotable.plot_ppt(self.shapes)

    def extract_swimlane_data(self):
        """
        After some thought am taking a very simple approach here.

        Just go through each activity on the plan and:
        - Check which swimlane it is for
        - Look at the track number and height in tracks and therefore work out the bottom track number
        - If the bottom track number is larger than that already recorded, update to reflect.
        - Once you have checked every activity, the entry in each swimlane will have recorded the number of tracks
          required for that swimlane.
        - We can then go back and calculate the start and end track for each swimlane which is what the plot method
          needs.  If a specific swimlane order is dictated then that is used here.  Where a swimlane doesn't appear in
          the ordering then we just add it to the end.

          So we will end up with a dict, with one entry for each (named) swimlane, and against each swimlane we will
          see the start track number and the end track number.

        :return:
        """

        swimlane_data = {}
        swimlane_manager = SwimlaneManager(self.swimlanes)
        for record in self.plan_data:
            swimlane = record.activity_layout_attributes.swimlane_name
            track_num_high = \
                record.activity_layout_attributes.track_number + \
                record.activity_layout_attributes.number_of_tracks_to_span \
                - 1
            if swimlane not in swimlane_data:
                # Adding record for this swimlane. Also remember ordering as is important later.
                swimlane_record = {
                    'swimlane_number': swimlane_manager.get_swimlane_number(swimlane),
                    'highest_track_within_lane': track_num_high
                }
                swimlane_data[swimlane] = swimlane_record
            else:
                swimlane_record = swimlane_data[swimlane]
                if track_num_high > swimlane_record['highest_track_within_lane']:
                    swimlane_record['highest_track_within_lane'] = track_num_high

        # We now have a dict containing a record for each swimlane of lowest and highest relative track number used.
        # Can now calculate the start and end track number for each swimlane - in order that lanes were encountered.

        swimlane_plot_data = {}
        end_track = 0

        # Create list of swimlane data entries in order of swimlane number
        sorted_entries = sorted(swimlane_data.items(), key=lambda x: x[1]['swimlane_number'])

        # Calculate start and end track number for each swimlane based on number of tracks required for swimlane
        for swimlane, swimlane_entry in sorted_entries:
            start_track = end_track + 1
            end_track = start_track + swimlane_entry['highest_track_within_lane'] - 1
            swimlane_plot_data[swimlane] = {
                'start_track': start_track,
                'end_track': end_track
            }
        return swimlane_plot_data

    def align_months(self):
        """
        If earliest or latest dates haven't been specified explicitly, then calculate from plan data.
        Then adjust to be first and last days of month respectively to ensure that month bar
        aligns with configured plot area.

        :return:
        """

        if self.plot_driver.min_start_date is None:
            start_dates = [record.start_date for record in self.plan_data]
            self.plot_driver.min_start_date = reduce(
                lambda min_date, start_date: start_date if start_date < min_date else min_date, start_dates)

        if self.plot_driver.max_end_date is None:
            end_dates = [record.end_date for record in self.plan_data]
            self.plot_driver.max_end_date = reduce(
                lambda max_date, end_date: end_date if not end_date is None and end_date > max_date else max_date,
                end_dates)

        # Regardless of whether start and end dates have been configured, we need to align with whole month
        self.plot_driver.min_start_date = first_day_of_month(self.plot_driver.min_start_date)
        self.plot_driver.max_end_date = last_day_of_month(self.plot_driver.max_end_date)

    def plot_vertical_line(self, current_date):
        x = self.plot_driver.date_to_x_coordinate(current_date, "start")
        top = self.plot_config.top
        bottom = self.plot_config.bottom
        line = self.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x, top, x, bottom)
        today_line_format = self.format_config['today_line']
        today_line_colour = today_line_format['line_rgb']
        line.line.color.rgb = RGBColor(*today_line_colour)
