from pptx.dml.color import RGBColor
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Cm
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from source.visualiser.shape_formatting import ShapeFormatting
from source.visualiser.text_formatting import TextFormatting
from source.visualiser.visual_element_shape import VisualElementShape


class PlotableElement:
    """
    Represents the data required to plot a shape within the visual.

    The intention is that the data can be used to plot in any medium, although the initial use case
    will be to plot on a PowerPoint slide.

    The information stored includes:

    Always Include:
    - Shape to plot
    - top
    - left
    - bottom
    - right

    Optional Depending Upon Case:
    - text

    Shape Formatting Information:
    - Fill colour etc.

    Text Formatting Information:
    - Horizontal alignment etc.

    Derived Values:
    - width
    - height

    Note that depending upon case the user may required different ways of expressing
    """
    def __init__(
            self,
            shape: VisualElementShape,
            top: Cm,
            left: Cm,
            bottom: Cm,
            right: Cm,
            shape_formatting: ShapeFormatting,
            text: str = None,
            text_formatting: TextFormatting = None
    ):
        """
        :param shape:
        :param top:
        :param left:
        :param bottom:
        :param right:
        :param shape_formatting:
        """
        self.shape = shape
        self.top = top
        self.left = left
        self.bottom = bottom
        self.right = right
        self.shape_formatting = shape_formatting
        self.text = text
        self.text_formatting = text_formatting

    @property
    def width(self):
        return self.right - self.left

    @property
    def height(self):
        return self.bottom - self.top

    def plot_ppt(self, shapes_object: SlideShapes):
        """
        Plots a shape on a PowerPoint slide given a shapes object.
        :return:
        """
        plotted_shape = shapes_object.add_shape(
            self.shape.ppt_shape,
            round(self.left),
            round(self.top),
            round(self.width),
            round(self.height)
        )

        # Adjust rounded corner radius, but only if the shape has corners
        if self.shape == VisualElementShape.ROUNDED_RECTANGLE:
            target_radius = self.shape_formatting.corner_radius
            adjustment_value = target_radius / self.height
            plotted_shape.adjustments[0] = adjustment_value

        fill = plotted_shape.fill
        line = plotted_shape.line

        if self.shape_formatting.fill_colour is None:
            fill.background()
            line.fill.background()
        else:
            fill.solid()

            fill_colour_ppt = self.rgb_ppt_format(self.shape_formatting.fill_colour.get_rgb())
            fill.fore_color.rgb = RGBColor(*fill_colour_ppt)

            line.color.rgb = RGBColor(*self.rgb_ppt_format(self.shape_formatting.line_colour.get_rgb()))

        # If the object has text, then add to the shape
        if self.text is not None:
            text_frame = plotted_shape.text_frame

            # Adjust text margin depending upon positioning. To help readability by having small gap

            text_frame.margin_top = self.text_formatting.margin_top
            text_frame.margin_bottom = self.text_formatting.margin_bottom
            text_frame.margin_left = self.text_formatting.margin_left
            text_frame.margin_right = self.text_formatting.margin_right
            text_frame.vertical_anchor = self._text_vertical_alignment(self.text_formatting.vertical_align)

            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = self.text

            font = run.font
            paragraph.line_spacing = 0.8
            paragraph.alignment = self._text_horizontal_alignment(self.text_formatting.horizontal_align)

            font.name = 'Calibri'  # Hard-coded for now
            font.size = self.text_formatting.font_size
            font.bold = self.text_formatting.font_bold
            font.italic = self.text_formatting.font_italic
            font.color.rgb = RGBColor(*self.rgb_ppt_format(self.text_formatting.font_colour.get_rgb()))
        return plotted_shape

    @staticmethod
    def _text_vertical_alignment(alignment):
        if alignment == "top":
            return MSO_ANCHOR.TOP
        elif alignment == "bottom":
            return MSO_ANCHOR.BOTTOM
        else:
            return MSO_ANCHOR.MIDDLE

    @staticmethod
    def _text_horizontal_alignment(format_text_align):
        """
        Takes the alignment field from the element formatting data and converts to the appropriate value for the
        pptx-python setting.

        :param format_text_align:
        :return:
        """
        if format_text_align == "left":
            return PP_ALIGN.LEFT
        elif format_text_align == "right":
            return PP_ALIGN.RIGHT
        elif format_text_align == "centre":
            return PP_ALIGN.CENTER
        else:
            # Default to centre
            return PP_ALIGN.CENTER

    @staticmethod
    def rgb_ppt_format(rgb):
        return map(lambda x: int(x*255), rgb)

