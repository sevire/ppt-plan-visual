import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Cm

from source.excel_plan import ExcelPlan
from source.plot_driver import PlotDriver
from source.tests.test_data.test_data_01 import plot_area_config
from source.utilities import get_path_name_ext


class PlanVisualiser:
    """
    Object which manages the creation of a PPT slide with a visual representation
    of a project plan (or similar).

    The supplied data

    :param plan_data:
    """

    def __init__(self, plan_data, plot_config, format_config, template_path):
        # The actual plan data with activities and milestones, start/finish dates etc.
        self.plan_data = plan_data

        # Data to define plot area for elements, tracks etc.
        self.plot_config = plot_config

        # Data with pre-determined formatting properties to apply to elements.
        self.format_config = format_config

        self.template = template_path

        folder, base, ext = get_path_name_ext(template_path)
        self.slides_out_path = os.path.join(folder, base + '_out' + ext)
        self.prs = Presentation(template_path)

        visual_slide = self.prs.slides[0]  # Assume there is one slide and that's where we will place the visual

        self.shapes = visual_slide.shapes
        self.plot_driver = PlotDriver(plot_area_config)

        self.swimlane_driver = self.extract_swimlane_data()

    @classmethod
    def from_excel(cls, plan_data_excel_path, excel_driver_config, template_path):
        excel_manager = ExcelPlan(excel_driver_config, plan_data_excel_path)

        extracted_plot_config = excel_manager.extract_plot_config_data()
        extracted_format_config = excel_manager.extract_format_config_data()
        extracted_plan_data = excel_manager.read_plan_data()

        return PlanVisualiser(extracted_plan_data, extracted_plot_config, extracted_format_config, template_path)

    def extract_swimlane_data(self):
        """
        After some thought am taking a very simple approach here.

        Just go through each activity on the plan and:
        - Check which swimlane it is for
        - Look at the track number and height in tracks and therefore work out the bottom track number
        - If the bottom track number is larger than that already recorded, update to reflect.
        - Once you have checked every activity, the entry in each swimlane will have recorded the number of tracks
          required for that swimlane.
        - We can then go back and calculate the start and end track for each swimlane which is what the plot method
          needs.

          So we will end up with a dict, with one entry for each (named) swimlane, and against each swimlane we will
          see the start track number and the end track number.

        :return:
        """

        swimlane_data = {}
        entry_num = 0
        for record in self.plan_data:
            swimlane = record['swimlane']
            track_num_high = record['track_num'] + record['bar_height_in_tracks'] - 1
            if swimlane not in swimlane_data:
                # Adding record for this swimlane. Also remember ordering as is important later.
                swimlane_record = {
                    'swimlane_order': entry_num + 1,
                    'highest_track_within_lane': track_num_high
                }
                swimlane_data[swimlane] = swimlane_record
                entry_num += 1
            else:
                swimlane_record = swimlane_data[swimlane]
                if track_num_high > swimlane_record['highest_track_within_lane']:
                    swimlane_record['highest_track_within_lane'] = track_num_high

        # We now have a dict containing a record for each swimlane of lowest and highest relative track number used.
        # Can now calculate the start and end track number for each swimlane - in order that lanes were encountered.

        swimlane_plot_data = {}
        end_track = 0
        for lane_number in range(0, len(swimlane_data)):
            swimlane_entries = [(key, swimlane_data[key]) for key in swimlane_data.keys() if swimlane_data[key]['swimlane_order'] == lane_number + 1]

            assert(len(swimlane_entries) == 1)

            swimlane, swimlane_entry = swimlane_entries[0]
            start_track = end_track + 1
            end_track = end_track + swimlane_entry['highest_track_within_lane'] - 1
            swimlane_plot_data[swimlane] = {
                'start_track': start_track,
                'end_track': end_track
            }
        return swimlane_plot_data

    def plot_slide(self):
        """
        Opens a supplied template file in order to allow consistency with other slides in a deck.

        Then creates a new slide with a plotted plan visual based on the data read in to the object.

        Then writes the one-slide deck to a different filename in the same folder.

        :return:
        """

        for plotable_element in self.plan_data:
            start = plotable_element['start_date']
            end = plotable_element['end_date']
            swimlane = plotable_element['swimlane']
            track_num = plotable_element['track_num']
            num_tracks = plotable_element['bar_height_in_tracks']
            shape_format = plotable_element['format_properties']
            if plotable_element['type'] == 'bar':
                self.plot_bar(start, end, swimlane, track_num, num_tracks, shape_format)
            elif plotable_element['type'] == 'milestone':
                self.plot_milestone(start, swimlane, track_num, shape_format)

        self.prs.save(self.slides_out_path)

    def plot_bar(self, start_date, end_date, swimlane, track_number, num_tracks, format_properties):
        swimlane_start = self.swimlane_driver[swimlane]['start_track']

        left = self.plot_driver.date_to_x_coordinate(start_date)
        right = self.plot_driver.date_to_x_coordinate(end_date)

        top = self.plot_driver.track_number_to_y_coordinate(swimlane_start + track_number - 1) - 1
        height = self.plot_driver.height_of_track(num_tracks)

        shape = self.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, right-left, height
        )

        properties = self.format_config['format_categories'][format_properties]

        # Adjust rounded corner radius
        target_radius = properties['corner_radius']
        adjustment_value = target_radius / height

        shape.adjustments[0] = adjustment_value

        self.shape_fill(shape, properties)
        self.shape_line(shape, properties)

    def plot_milestone(self, start_date, swimlane, track_number, format_properties):
        swimlane_start = self.swimlane_driver[swimlane]['start_track']
        milestone_width = self.plot_config['milestone_width']
        milestone_height = self.plot_config['track_height']

        left = self.plot_driver.date_to_x_coordinate(start_date) - milestone_width / 2
        top = self.plot_driver.track_number_to_y_coordinate(swimlane_start + track_number - 1)

        shape = self.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, left, top, milestone_width, milestone_height
        )

        self.shape_fill(shape, self.format_config['format_categories'][format_properties])
        self.shape_line(shape, self.format_config['format_categories'][format_properties])

    def shape_fill(self, shape, format_properties):
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*format_properties['fill_rgb'])

    def shape_line(self, shape, format_properties):
        line = shape.line
        line.color.rgb = RGBColor(*format_properties['line_rgb'])




